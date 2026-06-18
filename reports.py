import pandas as pd
import tempfile

from pptx import Presentation
from pptx.util import Inches


def create_excel_report(df):

    temp_file = tempfile.NamedTemporaryFile(
        delete=False,
        suffix=".xlsx"
    )

    path = temp_file.name

    with pd.ExcelWriter(
        path,
        engine="xlsxwriter"
    ) as writer:

        df.to_excel(
            writer,
            sheet_name="Data",
            index=False
        )

        summary = pd.DataFrame({
            "Column": df.columns,
            "Datatype": df.dtypes.astype(str),
            "Missing Values": df.isnull().sum()
        })

        summary.to_excel(
            writer,
            sheet_name="Summary",
            index=False
        )

    return path


def create_ppt_report(
    insights,
    image_paths=None
):

    prs = Presentation()

    slide = prs.slides.add_slide(
        prs.slide_layouts[0]
    )

    slide.shapes.title.text = (
        "AI CSVision Report"
    )

    slide.placeholders[1].text = (
        "Generated automatically"
    )

    slide = prs.slides.add_slide(
        prs.slide_layouts[1]
    )

    slide.shapes.title.text = (
        "AI Insights"
    )

    text_frame = (
        slide.placeholders[1]
        .text_frame
    )

    for item in insights:

        p = text_frame.add_paragraph()

        p.text = item

    if image_paths:

        for img in image_paths:

            slide = prs.slides.add_slide(
                prs.slide_layouts[5]
            )

            slide.shapes.title.text = (
                "Visualization"
            )

            slide.shapes.add_picture(
                img,
                Inches(1),
                Inches(1),
                width=Inches(7)
            )

    temp_file = tempfile.NamedTemporaryFile(
        delete=False,
        suffix=".pptx"
    )

    prs.save(temp_file.name)

    return temp_file.name